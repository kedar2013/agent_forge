"""Tests for mcp_servers/slide_reporting_server.py.

Deliberately does NOT import anything from `app.*` (the Postgres-backed
agent_forge app) or use conftest.py's fixtures -- this module is fully
standalone (MySQL + local rendering only), and mixing `app.X` bare imports
with `backend.app.X` prefixed imports (as tests/conftest.py uses) in one
process causes a SQLAlchemy "Table already defined for this MetaData
instance" clash, which is a pre-existing issue in this repo unrelated to
slide_reporting_server.py (reproducible today on main via
`pytest --collect-only`, tests/test_agent_runtime.py). Keeping this file's
imports scoped to mcp_servers/ sidesteps it entirely.

Run from backend/:
    ../.venv/Scripts/python.exe -m pytest tests/test_slide_reporting.py -v
"""

import json
import os
import re
import sys

import pandas as pd
import pytest
from pptx import Presentation

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "mcp_servers"))

import slide_reporting_server as srs  # noqa: E402
from _db import get_connection  # noqa: E402


def _sales_db_available() -> bool:
    try:
        conn = get_connection()
        conn.close()
        return True
    except Exception:
        return False


DB_AVAILABLE = _sales_db_available()
requires_db = pytest.mark.skipif(not DB_AVAILABLE, reason="sales_analytics MySQL DB not reachable")


# ---------------------------------------------------------------------------
# SQL guardrails
# ---------------------------------------------------------------------------


@pytest.mark.parametrize(
    "sql",
    [
        "SELECT 1",
        "SELECT region_name FROM regions",
        "WITH x AS (SELECT 1 AS a) SELECT a FROM x",
        "select * from orders where order_status = 'Completed'",
    ],
)
def test_validate_single_select_accepts_valid_select(sql):
    assert srs._validate_single_select(sql) is None


@pytest.mark.parametrize(
    "sql",
    [
        "DELETE FROM orders",
        "DROP TABLE orders",
        "UPDATE orders SET order_status = 'Cancelled'",
        "SELECT 1; DROP TABLE orders",
        "CREATE TABLE evil (id INT)",
        "not even sql (((",
    ],
)
def test_validate_single_select_rejects_unsafe_sql(sql):
    assert srs._validate_single_select(sql) is not None


def test_ensure_limit_appends_when_missing():
    limited_sql, had_own_limit = srs._ensure_limit("SELECT * FROM orders", 25)
    assert had_own_limit is False
    assert "LIMIT 25" in limited_sql


def test_ensure_limit_preserves_existing_limit():
    limited_sql, had_own_limit = srs._ensure_limit("SELECT * FROM orders LIMIT 5", 25)
    assert had_own_limit is True
    assert "LIMIT 5" in limited_sql
    assert "LIMIT 25" not in limited_sql


# ---------------------------------------------------------------------------
# Chart-shape classification (rule-based, no DB/LLM)
# ---------------------------------------------------------------------------


def test_classify_shape_kpi_for_single_row():
    df = pd.DataFrame([{"total_revenue": 12345.67}])
    assert srs._classify_shape(df) == "kpi"


def test_classify_shape_bar_for_one_category_one_metric():
    df = pd.DataFrame({"region_name": ["East", "West", "North"], "revenue": [100, 200, 150]})
    assert srs._classify_shape(df) == "bar"


def test_classify_shape_table_when_too_many_categories():
    df = pd.DataFrame({"customer_name": [f"Customer {i}" for i in range(20)], "revenue": list(range(20))})
    assert srs._classify_shape(df) == "table"


def test_classify_shape_line_for_date_plus_metric():
    df = pd.DataFrame({"month": ["2024-01", "2024-02", "2024-03"], "revenue": [100, 110, 90]})
    assert srs._classify_shape(df) == "line"


def test_classify_shape_table_for_many_numeric_columns():
    df = pd.DataFrame({"a": [1], "b": [2], "c": [3], "d": [4], "e": [5]})
    assert srs._classify_shape(df) == "table"


# ---------------------------------------------------------------------------
# slide_builder_tool -- renders a synthetic plan, no DB/LLM needed
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_slide_builder_tool_renders_and_reloads(tmp_path):
    plan = srs.SlidePlan(
        title="Unit Test Deck",
        slides=[
            srs.SlideSpec(kind="title", heading="Unit Test Deck"),
            srs.SlideSpec(kind="kpi", heading="Total Revenue", chart_config={"value": "1.2M", "label": "revenue"}),
            srs.SlideSpec(
                kind="bar", heading="Revenue by Region",
                chart_config={"categories": ["East", "West"], "values": [100.0, 200.0], "y_label": "revenue"},
            ),
            srs.SlideSpec(
                kind="line", heading="Monthly Trend",
                chart_config={"categories": ["2024-01", "2024-02"], "values": [100.0, 120.0], "y_label": "revenue"},
            ),
            srs.SlideSpec(
                kind="table", heading="All Customers",
                table_data=[{"customer_name": f"Customer {i}", "revenue": i} for i in range(15)],
            ),
            srs.SlideSpec(kind="bullets", heading="Key Takeaways", bullets=["Revenue grew 20%.", "East led all regions."]),
        ],
    )

    msg = await srs.slide_builder_tool(plan.model_dump_json(), "Unit test question")
    assert msg.startswith("Presentation generated:")

    filename = msg.rsplit("/", 1)[-1].split("?", 1)[0]
    path = os.path.join(srs.OUTPUT_DIR, filename)
    try:
        assert os.path.exists(path)
        prs = Presentation(path)
        slide_count = len(prs.slides._sldIdLst)
        # title + kpi + bar + line + 2 table pages (15 rows / 12-per-slide) + bullets
        assert slide_count == 7
    finally:
        if os.path.exists(path):
            os.remove(path)


@pytest.mark.asyncio
async def test_slide_builder_tool_empty_result_produces_no_data_slide():
    plan = srs.SlidePlan(
        title="Empty Result",
        slides=[
            srs.SlideSpec(kind="title", heading="Empty Result"),
            srs.SlideSpec(kind="bullets", heading="No Data", bullets=["No data found for this query."]),
        ],
    )
    msg = await srs.slide_builder_tool(plan.model_dump_json(), "A question with no matching rows")
    assert msg.startswith("Presentation generated:")

    filename = msg.rsplit("/", 1)[-1].split("?", 1)[0]
    path = os.path.join(srs.OUTPUT_DIR, filename)
    try:
        prs = Presentation(path)
        assert len(prs.slides._sldIdLst) == 2
    finally:
        if os.path.exists(path):
            os.remove(path)


@pytest.mark.asyncio
async def test_slide_builder_tool_invalid_plan_returns_error_not_exception():
    msg = await srs.slide_builder_tool("not json", "question")
    assert msg.startswith("Couldn't build the presentation:")


# ---------------------------------------------------------------------------
# Integration: the 3 example questions end-to-end against the live,
# already-seeded sales_analytics DB (skipped if it isn't reachable). These
# also make real Gemini calls (nl_to_sql_tool / chart_planner_tool's bullets).
# ---------------------------------------------------------------------------


def _generated_files_to_clean(msg: str) -> str | None:
    match = re.search(r"/generated-files/([^\s]+\.pptx)", msg)
    return match.group(1) if match else None


@requires_db
@pytest.mark.asyncio
async def test_example_question_revenue_by_region_2025():
    question = "Show total revenue by region for 2025"
    sql_result = await srs.nl_to_sql_tool(question)
    sql_dict = json.loads(sql_result)
    assert "error" not in sql_dict, sql_dict.get("error")

    exec_result = await srs.sql_execution_tool(sql_dict["sql"])
    exec_dict = json.loads(exec_result)
    assert "error" not in exec_dict, exec_dict.get("error")

    plan_result = await srs.chart_planner_tool(exec_result, question)
    plan_dict = json.loads(plan_result)
    assert "error" not in plan_dict, plan_dict.get("error")
    kinds = [s["kind"] for s in plan_dict["slides"]]
    assert "bar" in kinds or "table" in kinds

    deck_msg = await srs.slide_builder_tool(plan_result, question)
    assert deck_msg.startswith("Presentation generated:")
    filename = _generated_files_to_clean(deck_msg)
    path = os.path.join(srs.OUTPUT_DIR, filename)
    try:
        prs = Presentation(path)
        assert len(prs.slides._sldIdLst) > 0
    finally:
        if os.path.exists(path):
            os.remove(path)


@requires_db
@pytest.mark.asyncio
async def test_example_question_monthly_trend_2024():
    question = "What's the monthly revenue trend for 2024?"
    sql_result = await srs.nl_to_sql_tool(question)
    sql_dict = json.loads(sql_result)
    assert "error" not in sql_dict, sql_dict.get("error")

    exec_result = await srs.sql_execution_tool(sql_dict["sql"])
    exec_dict = json.loads(exec_result)
    assert "error" not in exec_dict, exec_dict.get("error")

    plan_result = await srs.chart_planner_tool(exec_result, question)
    plan_dict = json.loads(plan_result)
    assert "error" not in plan_dict, plan_dict.get("error")
    kinds = [s["kind"] for s in plan_dict["slides"]]
    assert "line" in kinds or "table" in kinds

    deck_msg = await srs.slide_builder_tool(plan_result, question)
    assert deck_msg.startswith("Presentation generated:")
    filename = _generated_files_to_clean(deck_msg)
    path = os.path.join(srs.OUTPUT_DIR, filename)
    try:
        prs = Presentation(path)
        assert len(prs.slides._sldIdLst) > 0
    finally:
        if os.path.exists(path):
            os.remove(path)


@requires_db
@pytest.mark.asyncio
async def test_example_question_top_10_customers():
    question = "List the top 10 customers by total order value"
    sql_result = await srs.nl_to_sql_tool(question)
    sql_dict = json.loads(sql_result)
    assert "error" not in sql_dict, sql_dict.get("error")

    exec_result = await srs.sql_execution_tool(sql_dict["sql"])
    exec_dict = json.loads(exec_result)
    assert "error" not in exec_dict, exec_dict.get("error")
    assert exec_dict["row_count"] <= 10

    plan_result = await srs.chart_planner_tool(exec_result, question)
    plan_dict = json.loads(plan_result)
    assert "error" not in plan_dict, plan_dict.get("error")

    deck_msg = await srs.slide_builder_tool(plan_result, question)
    assert deck_msg.startswith("Presentation generated:")
    filename = _generated_files_to_clean(deck_msg)
    path = os.path.join(srs.OUTPUT_DIR, filename)
    try:
        prs = Presentation(path)
        assert len(prs.slides._sldIdLst) > 0
    finally:
        if os.path.exists(path):
            os.remove(path)
