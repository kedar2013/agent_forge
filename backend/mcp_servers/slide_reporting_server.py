"""A real MCP server that turns a natural-language sales question into a
downloadable PPTX deck: NL->SQL (direct Gemini call) -> guarded MySQL
execution against `sales_analytics` -> rule-based chart/slide planning ->
python-pptx + matplotlib rendering. Runs over stdio, spawned as a subprocess
by Agent Forge's mcp_tool (StdioConnectionParams), following the same pattern
as the other mcp_servers/*.py files.

Exposes FOUR separate tools rather than one coarse tool, on purpose: the
calling agent (slide_reporting_agent, see
scripts/seed_slide_reporting_agent.py) is instructed to call them in strict
order -- nl_to_sql_tool -> sql_execution_tool -> chart_planner_tool ->
slide_builder_tool -- piping each JSON result's relevant fields into the
next call's arguments, with a one-retry self-correction loop back to
nl_to_sql_tool if sql_execution_tool errors.

SQL safety is layered:
  1. sqlglot parses the SQL into a real AST (dialect="mysql") -- a query is
     rejected unless it's exactly one `SELECT` (or `WITH ... SELECT`)
     statement, structurally, not by keyword/string matching. Any
     INSERT/UPDATE/DELETE/DROP/ALTER/CREATE/TRUNCATE node anywhere in the
     tree (not just top-level) is also rejected.
  2. Both nl_to_sql_tool (right after generation) and sql_execution_tool
     (independently, in case it's ever called with hand-written SQL)
     re-validate before touching the database.
  3. A LIMIT is appended via an AST edit (`Select.limit(...)`), never string
     concatenation, only when the query doesn't already have one.
  4. Defense-in-depth beyond this file: point SALES_DB_USER at a MySQL user
     granted SELECT-only on sales_analytics
     (`GRANT SELECT ON sales_analytics.* TO 'reporting_agent'@'%'`) rather
     than a superuser, so even a bug here still can't write.

Files are written to backend/generated_files/ (same directory
document_export_server.py uses, already mounted by app/main.py at
/generated-files) with an ABSOLUTE download URL via BACKEND_PUBLIC_URL.

Run standalone for a smoke test:
    python mcp_servers/slide_reporting_server.py
"""

import json
import os
import re
import time
import uuid
from datetime import datetime
from decimal import Decimal
from typing import Any, Literal

import pandas as pd
import sqlglot
from mcp.server.fastmcp import FastMCP
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pydantic import BaseModel
from sqlglot import exp

from _chart_utils import fmt_number as _fmt_number, render_bar_chart, render_line_chart
from _signed_urls import sign_filename
from _db import get_connection
from _watermark import watermark_pptx_slide
from _schema_context import SCHEMA_CONTEXT
from _theme import THEME

mcp = FastMCP("slide-reporting")

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "..", "generated_files")
TEMP_CHART_DIR = os.path.join(OUTPUT_DIR, "_chart_tmp")
PUBLIC_BASE_URL = os.environ.get("BACKEND_PUBLIC_URL", "http://127.0.0.1:8000")
GEMINI_MODEL = os.environ.get("GEMINI_MODEL", "gemini-2.5-flash")

_ROWS_PER_TABLE_SLIDE = 12


# ---------------------------------------------------------------------------
# Pydantic result models (the JSON contract each tool returns)
# ---------------------------------------------------------------------------


class NLToSQLResult(BaseModel):
    sql: str
    explanation: str
    expected_columns: list[str]
    confidence: Literal["high", "medium", "low"]


class SQLExecutionResult(BaseModel):
    row_count: int
    columns: list[str]
    data: list[dict[str, Any]]
    truncated: bool
    execution_ms: float


class SlideSpec(BaseModel):
    kind: Literal["title", "kpi", "bar", "line", "table", "bullets"]
    heading: str
    chart_config: dict[str, Any] | None = None
    table_data: list[dict[str, Any]] | None = None
    bullets: list[str] | None = None


class SlidePlan(BaseModel):
    title: str
    slides: list[SlideSpec]


# ---------------------------------------------------------------------------
# SQL guardrails -- shared by nl_to_sql_tool and sql_execution_tool
# ---------------------------------------------------------------------------

_FORBIDDEN_NODE_TYPES = (
    exp.Insert, exp.Update, exp.Delete, exp.Drop, exp.Alter, exp.Create,
    exp.TruncateTable,
)

# Defense-in-depth against function calls that read/write outside the
# database entirely -- a "reject INSERT/UPDATE/DDL node types" check does
# NOT catch these, since e.g. `SELECT LOAD_FILE('/etc/passwd')` is a
# syntactically ordinary SELECT with no write/DDL node anywhere in it. This
# tool has no per-call table allow-list to fall back on (it's a
# multi-table reporting tool by design, unlike app/tool_registry/
# data_query_tool.py's single-table sibling, which additionally requires
# the query reference its one configured table) -- this denylist is the
# one control that still applies regardless of which tables a query touches.
# Duplicated (not imported) from
# app/tool_registry/data_query_tool.py._FORBIDDEN_FUNCTION_NAMES: this
# script runs as its own subprocess with mcp_servers/ as sys.path[0] (see
# its sibling imports above, e.g. `from _db import get_connection`), not
# with the `app` package importable, so sharing the constant isn't safe.
_FORBIDDEN_FUNCTION_NAMES = frozenset(
    {
        "LOAD_FILE",  # MySQL: arbitrary local file read
        "SLEEP", "BENCHMARK",  # MySQL: timing-based DoS / blind enumeration
    }
)


def _forbidden_function_call(parsed: exp.Expression) -> str | None:
    for node in parsed.walk():
        candidate = node[0] if isinstance(node, tuple) else node
        if isinstance(candidate, exp.Func):
            name = (candidate.name or "").upper()
            if name in _FORBIDDEN_FUNCTION_NAMES:
                return f"Query calls a forbidden function ('{name}')."
    return None


def _validate_single_select(sql: str) -> str | None:
    """Returns an error message if `sql` isn't exactly one read-only SELECT
    statement, else None. Structural (via sqlglot's AST), not string
    matching -- see module docstring."""
    try:
        statements = [s for s in sqlglot.parse(sql, dialect="mysql") if s is not None]
    except Exception as exc:
        return f"SQL failed to parse: {exc}"

    if len(statements) != 1:
        return "Only a single SQL statement is allowed (no stacked statements)."

    stmt = statements[0]
    if not isinstance(stmt, exp.Select):
        return f"Only SELECT (or WITH ... SELECT) statements are allowed, got: {type(stmt).__name__}"

    for node in stmt.walk():
        if isinstance(node, _FORBIDDEN_NODE_TYPES):
            return "Query contains a forbidden write/DDL operation."

    return _forbidden_function_call(stmt)


def _ensure_limit(sql: str, max_rows: int) -> tuple[str, bool]:
    """Appends a LIMIT via an AST edit (never string concatenation) if the
    query doesn't already have one. Returns (sql_with_limit, had_own_limit)
    so the caller can tell a user-requested LIMIT (e.g. "top 10") apart from
    one we added ourselves for safety."""
    parsed = sqlglot.parse_one(sql, dialect="mysql")
    had_own_limit = parsed.args.get("limit") is not None
    if not had_own_limit:
        parsed = parsed.limit(max_rows)
    return parsed.sql(dialect="mysql"), had_own_limit


def _to_json_safe(value: Any) -> Any:
    if isinstance(value, Decimal):
        return float(value)
    if hasattr(value, "isoformat"):  # date/datetime/time
        return value.isoformat()
    return value


def _rows_to_json_safe(rows: list[dict]) -> list[dict]:
    return [{k: _to_json_safe(v) for k, v in row.items()} for row in rows]


# ---------------------------------------------------------------------------
# Formatting helpers
# ---------------------------------------------------------------------------


def _fmt_number_safe(value: Any) -> str:
    if isinstance(value, (int, float)) and not isinstance(value, bool):
        return _fmt_number(value)
    return str(value)


def _fmt_cell(value: Any) -> str:
    if value is None:
        return "—"
    if isinstance(value, bool):
        return str(value)
    if isinstance(value, float):
        return f"{value:,.2f}"
    if isinstance(value, int):
        return f"{value:,}"
    return str(value)


def _safe_filename(title: str) -> str:
    slug = re.sub(r"[^A-Za-z0-9]+", "-", title).strip("-").lower() or "presentation"
    return f"{slug[:60]}-{uuid.uuid4().hex[:8]}"


def _make_title(question: str) -> str:
    q = question.strip().rstrip("?").strip()
    return (q[:80] if q else "Sales Report") or "Sales Report"


def _hex(color: str) -> RGBColor:
    return RGBColor.from_string(color.lstrip("#"))


# ---------------------------------------------------------------------------
# Tool 1: nl_to_sql_tool
# ---------------------------------------------------------------------------

_NL_TO_SQL_SCHEMA = {
    "type": "object",
    "properties": {
        "sql": {"type": "string", "description": "A single MySQL SELECT statement."},
        "explanation": {"type": "string", "description": "One sentence explaining what the query computes."},
        "expected_columns": {"type": "array", "items": {"type": "string"}},
        "confidence": {"type": "string", "enum": ["high", "medium", "low"]},
    },
    "required": ["sql", "explanation", "expected_columns", "confidence"],
}


def _generate_sql(question: str, context: str, feedback: str | None) -> dict:
    from google import genai
    from google.genai import types

    client = genai.Client()
    prompt = f"{SCHEMA_CONTEXT}\n\nQuestion: {question}\n"
    if context:
        prompt += f"Additional context carried over from earlier in the conversation: {context}\n"
    if feedback:
        prompt += (
            f"\nYour previous attempt failed validation with this error:\n{feedback}\n"
            "Fix the SQL and try again.\n"
        )
    prompt += (
        "\nWrite exactly one MySQL SELECT statement (WITH ... SELECT is fine) "
        "that answers the question using only the tables/columns above. Never "
        "write INSERT/UPDATE/DELETE/DROP/ALTER/CREATE or more than one "
        "statement. Respond with the sql, a one-line explanation, the "
        "expected output column names, and your confidence."
    )

    response = client.models.generate_content(
        model=GEMINI_MODEL,
        contents=prompt,
        config=types.GenerateContentConfig(
            response_mime_type="application/json",
            response_schema=_NL_TO_SQL_SCHEMA,
            temperature=0.1,
        ),
    )
    return json.loads(response.text)


@mcp.tool()
async def nl_to_sql_tool(question: str, context: str = "") -> str:
    """Turn a natural-language sales/orders/revenue question into a single
    MySQL SELECT statement against the sales_analytics database. Always call
    this FIRST, before sql_execution_tool. Returns a JSON object (sql,
    explanation, expected_columns, confidence) -- pass its "sql" field
    straight into sql_execution_tool next. If sql_execution_tool later
    returns an error, call this tool exactly ONE more time with that error
    message appended to `context` before giving up.

    Args:
        question: The user's natural-language question, verbatim.
        context: Optional extra context -- prior-turn filters carried down
            from the orchestrator, or a previous SQL error to correct.
            Leave blank if none.
    """
    feedback = None
    for _ in range(2):
        try:
            parsed = _generate_sql(question, context, feedback)
        except Exception as exc:
            return json.dumps({"error": f"SQL generation failed: {exc}"})

        sql = parsed.get("sql", "")
        error = _validate_single_select(sql)
        if error is None:
            result = NLToSQLResult(
                sql=sql,
                explanation=parsed.get("explanation", ""),
                expected_columns=parsed.get("expected_columns", []),
                confidence=parsed.get("confidence", "medium"),
            )
            return result.model_dump_json()
        feedback = error

    return json.dumps({"error": f"Could not produce a valid single SELECT statement: {feedback}"})


# ---------------------------------------------------------------------------
# Tool 2: sql_execution_tool
# ---------------------------------------------------------------------------


@mcp.tool()
async def sql_execution_tool(sql: str, max_rows: int = 500) -> str:
    """Execute a single read-only MySQL SELECT statement (the "sql" field
    from nl_to_sql_tool) against sales_analytics and return the rows. Always
    call this SECOND, right after nl_to_sql_tool. Returns a JSON object
    (row_count, columns, data, truncated, execution_ms) on success, or
    {"error": "..."} on failure -- on error, call nl_to_sql_tool exactly ONE
    more time with the error as feedback, then give up with a clear message
    if it fails again.

    Args:
        sql: A single SELECT statement.
        max_rows: Row cap; a LIMIT is added automatically if the query
            doesn't already have one.
    """
    error = _validate_single_select(sql)
    if error:
        return json.dumps({"error": error})

    try:
        limited_sql, had_own_limit = _ensure_limit(sql, max_rows)
    except Exception as exc:
        return json.dumps({"error": f"Could not prepare query: {exc}"})

    start = time.perf_counter()
    conn = None
    try:
        conn = get_connection()
        with conn.cursor() as cursor:
            cursor.execute("SET SESSION MAX_EXECUTION_TIME=10000")
            cursor.execute(limited_sql)
            rows = cursor.fetchall()
    except Exception as exc:
        return json.dumps({"error": f"Query failed: {exc}"})
    finally:
        if conn is not None:
            conn.close()
    execution_ms = (time.perf_counter() - start) * 1000

    rows = _rows_to_json_safe(rows)
    columns = list(rows[0].keys()) if rows else []
    truncated = (not had_own_limit) and len(rows) >= max_rows

    result = SQLExecutionResult(
        row_count=len(rows), columns=columns, data=rows,
        truncated=truncated, execution_ms=round(execution_ms, 2),
    )
    return result.model_dump_json()


# ---------------------------------------------------------------------------
# Tool 3: chart_planner_tool
# ---------------------------------------------------------------------------


# Common key names, roughly in order of preference, that a data tool or an
# LLM-authored JSON blob might use for "the rows" -- deliberately broad since
# chart_planner_tool has to accept output from ANY upstream agent/tool (e.g.
# fund_analyst_agent's own free-form JSON), not just sql_execution_tool's/
# data_query_tool's shared {"data": [...]} contract.
_TABULAR_KEY_HINTS = (
    "data", "rows", "records", "results", "items", "entries", "holdings",
    "table", "table_data", "series", "list", "values",
)

_MISSING_TOKENS = {"", "-", "--", "n/a", "na", "null", "none", "nil"}

_FINANCE_SUFFIX_MULTIPLIERS = {
    "k": 1e3,
    "m": 1e6, "mm": 1e6, "mn": 1e6,
    "b": 1e9, "bn": 1e9,
    "cr": 1e7, "crore": 1e7, "crores": 1e7,
    "l": 1e5, "lac": 1e5, "lacs": 1e5, "lakh": 1e5, "lakhs": 1e5,
}
_FINANCE_NUMBER_RE = re.compile(r"^(-?\d+(?:\.\d+)?)\s*([a-zA-Z]+)?$")


def _parse_finance_number(raw: Any) -> float | None:
    """Parses a number out of the assorted string formats financial data
    shows up in -- "$1,234.56", "12.5%", "(340.2)" (accounting-style
    negative), "1.2M" / "3.4Cr" / "50 Lakh" (magnitude suffixes, including
    Indian-market ones), "N/A" / "-" (missing) -- returning None (never
    raising) for anything that isn't actually a number. Plain int/float pass
    straight through. This is what makes _coerce_numeric_columns generic
    enough for a fund analyst's formatting, not just sql_execution_tool's
    plain decimal strings."""
    if isinstance(raw, bool):
        return None
    if isinstance(raw, (int, float)):
        return float(raw)
    if not isinstance(raw, str):
        return None

    s = raw.strip()
    if s.lower() in _MISSING_TOKENS:
        return None

    negative = s.startswith("(") and s.endswith(")")
    if negative:
        s = s[1:-1].strip()

    s = s.replace(",", "").strip()
    s = re.sub(r"^[₹$€£]\s*", "", s)
    if s.endswith("%"):
        s = s[:-1].strip()

    match = _FINANCE_NUMBER_RE.match(s)
    if not match:
        return None
    number_part, suffix = match.groups()

    multiplier = 1.0
    if suffix:
        multiplier = _FINANCE_SUFFIX_MULTIPLIERS.get(suffix.lower())
        if multiplier is None:
            return None

    try:
        value = float(number_part) * multiplier
    except ValueError:
        return None
    return -value if negative else value


def _looks_tabular(value: Any) -> bool:
    return isinstance(value, list) and bool(value) and all(isinstance(v, dict) for v in value)


def _looks_columnar(value: Any) -> bool:
    """{"col_a": [...], "col_b": [...]} -- column-oriented rows, the other
    common tabular JSON shape besides row-oriented list-of-dicts."""
    if not isinstance(value, dict) or not value:
        return False
    lists = list(value.values())
    return all(isinstance(v, list) for v in lists) and len({len(v) for v in lists}) == 1


def _columnar_to_rows(value: dict) -> list[dict]:
    keys = list(value.keys())
    length = len(value[keys[0]]) if keys else 0
    return [{k: value[k][i] for k in keys} for i in range(length)]


def _find_tabular(value: Any, _depth: int = 0, _max_depth: int = 4) -> list[dict] | None:
    """Recursively hunts an arbitrary JSON value for the most plausible
    table, so chart_planner_tool isn't limited to data living under an exact
    "data" key at the top level -- any agent's own JSON, with its own key
    names and nesting, works too. Prefers a key from _TABULAR_KEY_HINTS at
    each level; otherwise falls back to the largest table found anywhere
    else in the structure."""
    if _looks_tabular(value):
        return value
    if _looks_columnar(value):
        return _columnar_to_rows(value)
    if _depth >= _max_depth:
        return None

    if isinstance(value, dict):
        for key in _TABULAR_KEY_HINTS:
            if key in value:
                found = _find_tabular(value[key], _depth + 1, _max_depth)
                if found:
                    return found
        candidates = [
            found for v in value.values()
            if (found := _find_tabular(v, _depth + 1, _max_depth))
        ]
        return max(candidates, key=len) if candidates else None

    if isinstance(value, list):
        candidates = [
            found for item in value
            if (found := _find_tabular(item, _depth + 1, _max_depth))
        ]
        return max(candidates, key=len) if candidates else None

    return None


def _extract_narrative(value: Any, _depth: int = 0, _max_depth: int = 3) -> list[str]:
    """Collects human-readable text out of a payload that has no tabular
    data at all (e.g. an agent that only ever answered in prose), so
    chart_planner_tool can still produce a bullets-only deck instead of
    erroring out."""
    if isinstance(value, str):
        text = value.strip()
        return [text] if text else []
    if _depth >= _max_depth:
        return []
    if isinstance(value, dict):
        bullets: list[str] = []
        for key in ("bullets", "summary", "narrative", "analysis", "insights", "text", "answer", "content", "message"):
            if key in value:
                found = _extract_narrative(value[key], _depth + 1, _max_depth)
                if found:
                    bullets.extend(found)
        return bullets
    if isinstance(value, list) and all(isinstance(v, str) for v in value):
        return [v.strip() for v in value if v.strip()]
    return []


def _normalize_execution_result(parsed: Any) -> dict:
    """Normalizes ANY upstream JSON payload -- not just the exact
    SQLExecutionResult shape sql_execution_tool/data_query_tool return --
    into {"data": [...]} (or {"error": "..."}), so the rest of
    chart_planner_tool only ever has to deal with one shape. Handles, in
    order:
      - the exact contract already ({"data": [...], ...}) -- used unchanged
      - one (or more) levels of accidental wrapping, e.g. an LLM re-presenting
        the prior tool's result under its own key instead of forwarding it
        verbatim (`{"query_facility_data_response": {...}}`)
      - {"columns": [...], "rows": [[...], ...]} -- column-list + row-matrix
      - column-oriented {"col": [...], "col2": [...]}
      - a bare list of dicts, list of lists, or list of scalars as the WHOLE
        payload (no wrapper object at all)
      - a flat single-record dict (no list anywhere) -- treated as one row
      - a JSON string containing any of the above (double-encoded payload)
    Falls back to {"data": [], "narrative": [...]} when there's no tabular
    data but there IS some text to summarize (e.g. a fund analyst agent that
    only answered in prose), and only returns {"error": ...} when truly
    nothing usable was found.
    """
    # Tolerate a JSON value that is itself a JSON-encoded string.
    for _ in range(2):
        if isinstance(parsed, str):
            try:
                parsed = json.loads(parsed)
            except Exception:
                break
        else:
            break

    if isinstance(parsed, dict) and _looks_tabular(parsed.get("data")):
        return parsed

    if isinstance(parsed, dict) and "error" in parsed and not parsed.get("data"):
        return {"error": parsed["error"]}

    if isinstance(parsed, dict) and isinstance(parsed.get("columns"), list) and isinstance(parsed.get("rows"), list):
        cols = parsed["columns"]
        rows = [dict(zip(cols, r)) for r in parsed["rows"] if isinstance(r, list)]
        if rows:
            return {"data": rows, "truncated": parsed.get("truncated", False)}

    if isinstance(parsed, dict) and isinstance(parsed.get("data"), dict):
        return {"data": [parsed["data"]]}

    if _looks_columnar(parsed):
        return {"data": _columnar_to_rows(parsed)}

    if _looks_tabular(parsed):
        return {"data": parsed}

    if isinstance(parsed, list) and parsed:
        if all(isinstance(v, list) for v in parsed):
            return {"data": [{f"col_{i}": c for i, c in enumerate(row)} for row in parsed]}
        if all(not isinstance(v, (dict, list)) for v in parsed):
            return {"data": [{"value": v} for v in parsed]}

    found = _find_tabular(parsed)
    if found:
        return {"data": found}

    if isinstance(parsed, dict) and parsed:
        scalar_items = {k: v for k, v in parsed.items() if not isinstance(v, (dict, list))}
        if scalar_items:
            return {"data": [scalar_items]}

    narrative = _extract_narrative(parsed)
    if narrative:
        return {"data": [], "narrative": narrative}

    if isinstance(parsed, dict):
        return {"error": "No tabular data could be found in the execution result."}
    return {"error": f"Could not make sense of the execution result (got {type(parsed).__name__})."}


def _coerce_numeric_columns(df: pd.DataFrame) -> pd.DataFrame:
    """Numeric values commonly arrive as JSON *strings* rather than JSON
    numbers -- plain decimals from a decimal-safe API like `data_query_tool`
    (JSON has no decimal type), or financial formatting like "$1,234.56",
    "12.5%", "(340.2)", "1.2M"/"3.4Cr" from an LLM-authored payload (e.g.
    fund_analyst_agent). Left alone, pandas gives those columns `object`
    dtype and every downstream `is_numeric_dtype` check in this file
    misclassifies them as categorical, which is what caused both the "No
    Data" and the `next()`-on-empty-generator crash seen in production.
    Coerces any column whose non-null values *all* parse as numbers (via
    _parse_finance_number) into a real numeric dtype; anything that doesn't
    parse cleanly (a genuine text column) is left untouched."""
    for col in df.columns:
        if pd.api.types.is_numeric_dtype(df[col]):
            continue
        non_null = df[col].dropna()
        if non_null.empty:
            continue
        parsed = non_null.map(_parse_finance_number)
        if parsed.notna().sum() == len(non_null):
            df[col] = df[col].map(lambda v: _parse_finance_number(v) if pd.notna(v) else v)
    return df


def _is_date_like(series: pd.Series) -> bool:
    if pd.api.types.is_datetime64_any_dtype(series):
        return True
    sample = series.dropna().astype(str).head(5)
    if sample.empty:
        return False
    try:
        pd.to_datetime(sample, errors="raise", format="mixed")
        return True
    except Exception:
        return False


def _classify_shape(df: pd.DataFrame) -> str:
    """Deterministic, rule-based shape classification (per spec thresholds)
    -- no LLM call in the common case. Returns "ambiguous" when nothing
    matches cleanly, so the caller can fall back to one small LLM call."""
    n_rows, n_cols = df.shape
    if n_rows == 1 and n_cols <= 2:
        return "kpi"

    numeric_cols = [c for c in df.columns if pd.api.types.is_numeric_dtype(df[c])]
    other_cols = [c for c in df.columns if c not in numeric_cols]
    date_cols = [c for c in other_cols if _is_date_like(df[c])]
    categorical_cols = [c for c in other_cols if c not in date_cols]

    if date_cols and numeric_cols:
        return "line"

    if len(numeric_cols) > 3:
        return "table"

    if len(categorical_cols) == 1 and len(numeric_cols) >= 1:
        return "table" if df[categorical_cols[0]].nunique() > 15 else "bar"

    if len(categorical_cols) == 2 and len(numeric_cols) == 1:
        # A grouped bar is possible here, but a table is the reliable
        # default -- clean grouped-bar rendering is real chart-writing
        # effort for uncertain benefit over a clear, always-correct table.
        return "table"

    return "ambiguous"


_SHAPE_SCHEMA = {
    "type": "object",
    "properties": {"shape": {"type": "string", "enum": ["kpi", "bar", "line", "table"]}},
    "required": ["shape"],
}


def _llm_classify_shape(df: pd.DataFrame, question: str) -> str:
    from google import genai
    from google.genai import types

    client = genai.Client()
    sample = df.head(5).to_dict("records")
    prompt = (
        f"Question: {question}\n"
        f"Result has {len(df)} rows and columns {list(df.columns)}. Sample rows: {sample}\n"
        "Pick the best slide visualization: 'kpi' for a single headline "
        "number, 'bar' for one category dimension vs one metric, 'line' for "
        "a date/time trend, or 'table' if it doesn't fit those cleanly."
    )
    try:
        response = client.models.generate_content(
            model=GEMINI_MODEL, contents=prompt,
            config=types.GenerateContentConfig(
                response_mime_type="application/json", response_schema=_SHAPE_SCHEMA, temperature=0.0
            ),
        )
        return json.loads(response.text).get("shape", "table")
    except Exception:
        return "table"


def _build_main_slide(shape: str, df: pd.DataFrame, question: str, execution_result: dict) -> SlideSpec:
    heading = _make_title(question)

    if shape == "kpi":
        value = df.iloc[0, -1]
        label = df.columns[-1] if len(df.columns) > 1 else (df.columns[0] if len(df.columns) else "Value")
        return SlideSpec(kind="kpi", heading=heading, chart_config={"value": _fmt_number_safe(value), "label": str(label)})

    numeric_cols = [c for c in df.columns if pd.api.types.is_numeric_dtype(df[c])]
    other_cols = [c for c in df.columns if c not in numeric_cols]

    # Both branches below need at least one column of the right kind — the
    # deterministic classifier already checked this, but shape can also come
    # from _llm_classify_shape's fuzzy judgment, which isn't guaranteed to
    # agree with what's actually in the (possibly just-coerced) dataframe.
    # Falling through to "table" on a mismatch beats crashing the whole tool
    # call (this exact gap previously surfaced as "coroutine raised
    # StopIteration" from a bare next() with no default).
    if shape == "bar" and other_cols and numeric_cols:
        cat_col = other_cols[0]
        val_col = numeric_cols[0]
        return SlideSpec(
            kind="bar", heading=heading,
            chart_config={
                "categories": df[cat_col].astype(str).tolist(),
                "values": [float(v) for v in df[val_col].tolist()],
                "y_label": str(val_col),
            },
        )

    if shape == "line" and numeric_cols:
        date_col = next((c for c in other_cols if _is_date_like(df[c])), None)
        if date_col is not None:
            val_col = numeric_cols[0]
            sorted_df = df.sort_values(date_col)
            return SlideSpec(
                kind="line", heading=heading,
                chart_config={
                    "categories": sorted_df[date_col].astype(str).tolist(),
                    "values": [float(v) for v in sorted_df[val_col].tolist()],
                    "y_label": str(val_col),
                },
            )

    table_heading = heading
    if execution_result.get("truncated"):
        table_heading += f" (showing first {execution_result.get('row_count')} rows)"
    return SlideSpec(kind="table", heading=table_heading, table_data=df.to_dict("records"))


def _compute_stats(df: pd.DataFrame) -> dict:
    stats: dict[str, Any] = {"row_count": len(df)}
    numeric_cols = [c for c in df.columns if pd.api.types.is_numeric_dtype(df[c])]
    if numeric_cols:
        main_col = numeric_cols[0]
        stats["metric"] = str(main_col)
        stats["max"] = float(df[main_col].max())
        stats["min"] = float(df[main_col].min())
        stats["total"] = float(df[main_col].sum())
        stats["top_row"] = {k: _fmt_number_safe(v) if isinstance(v, (int, float)) else str(v)
                             for k, v in df.loc[df[main_col].idxmax()].to_dict().items()}
    return stats


_INSIGHTS_SCHEMA = {
    "type": "object",
    "properties": {"bullets": {"type": "array", "items": {"type": "string"}}},
    "required": ["bullets"],
}


def _generate_insights(question: str, stats: dict) -> list[str]:
    from google import genai
    from google.genai import types

    client = genai.Client()
    prompt = (
        f"Question: {question}\n"
        f"Aggregated stats from the query result: {json.dumps(stats, default=str)}\n"
        "Write 2-4 short bullet-point takeaways (one sentence each) an "
        "executive would want on a summary slide. Be specific with the "
        "numbers given -- no filler, no restating the question."
    )
    try:
        response = client.models.generate_content(
            model=GEMINI_MODEL, contents=prompt,
            config=types.GenerateContentConfig(
                response_mime_type="application/json", response_schema=_INSIGHTS_SCHEMA, temperature=0.3
            ),
        )
        bullets = json.loads(response.text).get("bullets", [])
        return bullets or [f"Result contains {stats.get('row_count', 0)} rows."]
    except Exception:
        return [f"Result contains {stats.get('row_count', 0)} rows."]


@mcp.tool()
async def chart_planner_tool(execution_result_json: str, question: str) -> str:
    """Decide the chart type and slide outline from a data tool's or agent's
    result. Always call this THIRD, right after a successful
    sql_execution_tool call (or whatever tool/agent produced the data).
    Returns a JSON SlidePlan (title, slides: [...]) -- pass it straight into
    slide_builder_tool next.

    Accepts more than the exact SQLExecutionResult shape: a bare list of
    records, column-oriented data, a {"columns": [...], "rows": [[...]]}
    matrix, data nested one or more levels deep under any key name (e.g.
    "holdings", "results"), a single flat record, or even a JSON string
    that's itself double-encoded. If the payload has no tabular data at all
    but does have a text summary, it still produces a bullets-only deck
    instead of erroring out.

    Args:
        execution_result_json: The exact JSON string the data tool/agent you just
            called returned (sql_execution_tool, data_query_tool, a fund/market
            analyst agent's own JSON, or similar) — copy it verbatim. Do NOT
            re-type or summarize it; paste the tool's own output exactly as it
            came back to you.
        question: The user's original natural-language question.
    """
    try:
        parsed = json.loads(execution_result_json)
    except Exception as exc:
        return json.dumps({"error": f"Could not parse execution result: {exc}"})

    execution_result = _normalize_execution_result(parsed)

    if "error" in execution_result:
        return json.dumps({"error": execution_result["error"]})

    title = _make_title(question)
    rows = execution_result.get("data", [])
    slides = [SlideSpec(kind="title", heading=title)]

    if not rows:
        narrative = execution_result.get("narrative")
        if narrative:
            slides.append(SlideSpec(kind="bullets", heading="Summary", bullets=narrative[:8]))
        else:
            slides.append(SlideSpec(kind="bullets", heading="No Data", bullets=["No data found for this query."]))
        return SlidePlan(title=title, slides=slides).model_dump_json()

    df = _coerce_numeric_columns(pd.DataFrame(rows))
    shape = _classify_shape(df)
    if shape == "ambiguous":
        shape = _llm_classify_shape(df, question)

    slides.append(_build_main_slide(shape, df, question, execution_result))
    stats = _compute_stats(df)
    slides.append(SlideSpec(kind="bullets", heading="Key Takeaways", bullets=_generate_insights(question, stats)))

    return SlidePlan(title=title, slides=slides).model_dump_json()


# ---------------------------------------------------------------------------
# Tool 4: slide_builder_tool
# ---------------------------------------------------------------------------


def _add_heading(slide, text: str) -> None:
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = _hex(THEME["brand"])


def _add_title_slide(prs: Presentation, layout, title: str, question: str) -> None:
    slide = prs.slides.add_slide(layout)
    watermark_pptx_slide(slide, prs.slide_width, prs.slide_height)
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = _hex(THEME["brand"])

    sub_tb = slide.shapes.add_textbox(Inches(0.8), Inches(4.1), Inches(11.7), Inches(1.2))
    sub_tf = sub_tb.text_frame
    sub_tf.word_wrap = True
    q_p = sub_tf.paragraphs[0]
    q_p.alignment = PP_ALIGN.CENTER
    q_run = q_p.add_run()
    q_run.text = question
    q_run.font.size = Pt(15)
    q_run.font.color.rgb = _hex(THEME["ink_secondary"])

    ts_p = sub_tf.add_paragraph()
    ts_p.alignment = PP_ALIGN.CENTER
    ts_run = ts_p.add_run()
    ts_run.text = f"Generated on {datetime.now().strftime('%B %d, %Y %H:%M')}"
    ts_run.font.size = Pt(12)
    ts_run.font.italic = True
    ts_run.font.color.rgb = _hex(THEME["ink_muted"])


def _add_kpi_slide(prs: Presentation, layout, spec: SlideSpec) -> None:
    slide = prs.slides.add_slide(layout)
    watermark_pptx_slide(slide, prs.slide_width, prs.slide_height)
    _add_heading(slide, spec.heading)
    cfg = spec.chart_config or {}

    tb = slide.shapes.add_textbox(Inches(1), Inches(2.6), Inches(11.3), Inches(2))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(cfg.get("value", ""))
    run.font.size = Pt(72)
    run.font.bold = True
    run.font.color.rgb = _hex(THEME["sequential"])

    cap_tb = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.3), Inches(0.8))
    cap_p = cap_tb.text_frame.paragraphs[0]
    cap_p.alignment = PP_ALIGN.CENTER
    cap_run = cap_p.add_run()
    cap_run.text = str(cfg.get("label", ""))
    cap_run.font.size = Pt(18)
    cap_run.font.color.rgb = _hex(THEME["ink_secondary"])


def _add_chart_slide(prs: Presentation, layout, spec: SlideSpec, kind: str, temp_files: list[str]) -> None:
    slide = prs.slides.add_slide(layout)
    watermark_pptx_slide(slide, prs.slide_width, prs.slide_height)
    _add_heading(slide, spec.heading)
    cfg = spec.chart_config or {}
    categories = cfg.get("categories", [])
    values = cfg.get("values", [])
    y_label = cfg.get("y_label", "")

    tmp_path = os.path.join(TEMP_CHART_DIR, f"{uuid.uuid4().hex}.png")
    series = [{"name": y_label or "value", "values": values}]
    if kind == "bar":
        render_bar_chart(categories, series, tmp_path, y_label)
    else:
        render_line_chart(categories, series, tmp_path, y_label)
    temp_files.append(tmp_path)

    slide.shapes.add_picture(tmp_path, Inches(1.2), Inches(1.4), width=Inches(10.9))


def _add_bullets_slide(prs: Presentation, layout, spec: SlideSpec) -> None:
    slide = prs.slides.add_slide(layout)
    watermark_pptx_slide(slide, prs.slide_width, prs.slide_height)
    _add_heading(slide, spec.heading)
    tb = slide.shapes.add_textbox(Inches(1.2), Inches(1.6), Inches(10.9), Inches(5.2))
    tf = tb.text_frame
    tf.word_wrap = True
    bullets = spec.bullets or []
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        run = p.add_run()
        run.text = f"•  {bullet}"
        run.font.size = Pt(20)
        run.font.color.rgb = _hex(THEME["ink_primary"])


def _add_table_slides(prs: Presentation, layout, spec: SlideSpec) -> None:
    rows = spec.table_data or []
    if not rows:
        _add_bullets_slide(
            prs, layout,
            SlideSpec(kind="bullets", heading=spec.heading, bullets=["No data found for this query."]),
        )
        return

    columns = list(rows[0].keys())
    pages = [rows[i:i + _ROWS_PER_TABLE_SLIDE] for i in range(0, len(rows), _ROWS_PER_TABLE_SLIDE)]

    for page_idx, page_rows in enumerate(pages):
        slide = prs.slides.add_slide(layout)
        watermark_pptx_slide(slide, prs.slide_width, prs.slide_height)
        heading = spec.heading + (f" ({page_idx + 1}/{len(pages)})" if len(pages) > 1 else "")
        _add_heading(slide, heading)

        n_rows, n_cols = len(page_rows) + 1, len(columns)
        left, top, width = Inches(0.6), Inches(1.3), Inches(12.1)
        height = Inches(0.4 * n_rows)
        table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

        for c, col_name in enumerate(columns):
            cell = table.cell(0, c)
            cell.text = str(col_name)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _hex(THEME["brand"])
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        for r, row in enumerate(page_rows, start=1):
            for c, col_name in enumerate(columns):
                cell = table.cell(r, c)
                cell.text = _fmt_cell(row.get(col_name))
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(11)
                    p.font.color.rgb = _hex(THEME["ink_primary"])

        if len(rows) > _ROWS_PER_TABLE_SLIDE and page_idx == len(pages) - 1:
            note_tb = slide.shapes.add_textbox(left, top + height + Inches(0.15), width, Inches(0.4))
            note_run = note_tb.text_frame.paragraphs[0].add_run()
            note_run.text = f"Showing all {len(rows)} rows across {len(pages)} slides."
            note_run.font.size = Pt(11)
            note_run.font.italic = True
            note_run.font.color.rgb = _hex(THEME["ink_muted"])


@mcp.tool()
async def slide_builder_tool(slide_plan_json: str, question: str) -> str:
    """Render a SlidePlan (the exact JSON chart_planner_tool returned) into a
    downloadable PowerPoint deck. Always call this FOURTH and LAST. Returns a
    plain-text confirmation containing the download link -- relay it to the
    user along with a short natural-language summary of what the deck shows.

    Args:
        slide_plan_json: The exact JSON string chart_planner_tool returned.
        question: The user's original natural-language question.
    """
    try:
        plan = SlidePlan.model_validate(json.loads(slide_plan_json))
    except Exception as exc:
        return f"Couldn't build the presentation: invalid slide plan ({exc})"

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    os.makedirs(TEMP_CHART_DIR, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    temp_files: list[str] = []
    filename: str | None = None
    try:
        for spec in plan.slides:
            if spec.kind == "title":
                _add_title_slide(prs, blank_layout, spec.heading, question)
            elif spec.kind == "kpi":
                _add_kpi_slide(prs, blank_layout, spec)
            elif spec.kind in ("bar", "line"):
                _add_chart_slide(prs, blank_layout, spec, spec.kind, temp_files)
            elif spec.kind == "table":
                _add_table_slides(prs, blank_layout, spec)
            elif spec.kind == "bullets":
                _add_bullets_slide(prs, blank_layout, spec)

        if not plan.slides:
            return "Couldn't build the presentation: the slide plan was empty."

        filename = f"{_safe_filename(plan.title)}.pptx"
        path = os.path.join(OUTPUT_DIR, filename)
        prs.save(path)
    except Exception as exc:
        return f"Couldn't build the presentation: {exc}"
    finally:
        for f in temp_files:
            try:
                os.remove(f)
            except OSError:
                pass

    token = sign_filename("generated-files", filename)
    return f"Presentation generated: {PUBLIC_BASE_URL}/generated-files/{filename}?token={token}"


if __name__ == "__main__":
    mcp.run(transport="stdio")
